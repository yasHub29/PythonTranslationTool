# modules/pptx_translator/pptx_writer.py
# Replace text in-place for shapes identified by path tuples, preserve formatting and layout.

from pptx import Presentation

def _get_shape_by_path(slide, path):
    """
    Resolve a nested shape by path tuple.
    Example: path (2,0,3) => slide.shapes[2].shapes[0].shapes[3]
    """
    if not path:
        return None
    shape = slide.shapes[path[0]]
    for idx in path[1:]:
        # some shapes may not have .shapes; guard for safety
        if not hasattr(shape, "shapes"):
            return None
        shape = shape.shapes[idx]
    return shape

def _replace_paragraph_text_preserve_format(para, new_text):
    """
    Replace paragraph text while preserving the formatting of the first run.
    Strategy:
      - If paragraph has runs: set runs[0].text = new_text, clear others.
      - If no runs: set para.text = new_text
    """
    runs = para.runs
    if runs:
        # clear other runs but preserve formatting on runs[0]
        for i, r in enumerate(runs):
            if i == 0:
                r.text = new_text
            else:
                r.text = ""
    else:
        para.text = new_text

def write_pptx_from_template(src_path, dest_path, translated_slides):
    """
    src_path: original pptx (template)
    dest_path: destination file to save
    translated_slides: list aligned by slide index:
      {
        "translated_shape_texts": [ ((i1,i2,...), "translated text"), ... ],
        "images": [...]
      }
    """
    prs = Presentation(src_path)

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx >= len(translated_slides):
            break
        data = translated_slides[slide_idx]
        pairs = data.get("translated_shape_texts", [])

        for path, new_text in pairs:
            try:
                shape = _get_shape_by_path(slide, path)
                if shape is None:
                    continue
                # Only attempt if shape has a text_frame
                if not hasattr(shape, "text_frame") or shape.text_frame is None:
                    continue

                tf = shape.text_frame
                # Iterate paragraphs and replace text sequentially.
                # If the original paragraph count > 1, we preserve that structure.
                # We'll split new_text into lines and map line-wise to paragraphs.
                new_lines = new_text.splitlines() or [new_text]
                p_count = len(tf.paragraphs)
                line_index = 0

                for p in tf.paragraphs:
                    if line_index < len(new_lines):
                        _replace_paragraph_text_preserve_format(p, new_lines[line_index])
                        line_index += 1
                    else:
                        # No more translated lines: clear remaining paragraph text
                        _replace_paragraph_text_preserve_format(p, "")
                # If there are still extra lines (more translated lines than paragraphs),
                # append them to the last paragraph (preserve its formatting).
                if line_index < len(new_lines):
                    remaining = "\n".join(new_lines[line_index:])
                    last_para = tf.paragraphs[-1]
                    _replace_paragraph_text_preserve_format(last_para, last_para.text + ("\n" + remaining if last_para.text else remaining))

            except Exception:
                # If anything fails for a shape, skip it to avoid crashing translation for entire deck.
                continue

    prs.save(dest_path)

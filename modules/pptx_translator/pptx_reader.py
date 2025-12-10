# modules/pptx_translator/pptx_reader.py
# Extract structured text from PPTX: returns per-slide list of shape paths and texts.
# Works recursively for grouped shapes and keeps images list (unchanged).

from pptx import Presentation
from PIL import Image
from io import BytesIO

PICTURE_TYPE = 13

def _extract_from_shape(shape, path_prefix=()):
    """
    Return list of (path, text) for the given shape.
    path is a tuple of indices that identifies a nested shape inside the slide.
    """
    items = []
    # If shape has a text_frame and non-empty text, capture it
    if hasattr(shape, "text_frame") and shape.text_frame is not None:
        text = shape.text_frame.text or ""
        if text.strip():
            items.append((path_prefix, text))

    # If grouped shape, iterate children
    if hasattr(shape, "shapes"):
        for idx, child in enumerate(shape.shapes):
            items.extend(_extract_from_shape(child, path_prefix + (idx,)))
    return items

def read_pptx(path):
    """
    Returns:
      slides: [
        {
          'shape_texts': [ ((i1,i2,...), "original text"), ... ],
          'images': [PIL.Image, ...]
        }, ...
      ]
    The paths are relative to slide.shapes - to resolve a top-level shape, use index tuple (k,).
    """
    prs = Presentation(path)
    slides = []

    for slide in prs.slides:
        shape_texts = []
        images = []

        # iterate top-level shapes
        for top_index, shape in enumerate(slide.shapes):
            # collect text items from this top-level shape (including nested)
            items = _extract_from_shape(shape, (top_index,))
            shape_texts.extend(items)

            # images (keep as before)
            if getattr(shape, "shape_type", None) == PICTURE_TYPE:
                try:
                    img = shape.image
                    bio = BytesIO(img.blob)
                    pil = Image.open(bio)
                    images.append(pil.copy())
                except Exception:
                    pass

        slides.append({
            "shape_texts": shape_texts,
            "images": images
        })

    return slides
